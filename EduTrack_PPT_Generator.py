"""
EduTrack Professional PowerPoint Generator
==========================================
Generates a 25-slide client-ready presentation for the EduTrack School Management System.
Uses python-pptx library to create a polished, branded deck.

Usage:
    pip install python-pptx
    python EduTrack_PPT_Generator.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Brand Colors ──────────────────────────────────────────────────────────
BRAND_DARK    = RGBColor(0x0F, 0x17, 0x2A)   # Dark navy (#0F172A)
BRAND_BLUE    = RGBColor(0x2E, 0x5B, 0xFF)   # Primary blue (#2E5BFF)
BRAND_INDIGO  = RGBColor(0x63, 0x66, 0xF1)   # Indigo (#6366F1)
BRAND_PURPLE  = RGBColor(0x8B, 0x5C, 0xF6)   # Purple (#8B5CF6)
BRAND_TEAL    = RGBColor(0x14, 0xB8, 0xA6)   # Teal (#14B8A6)
BRAND_GREEN   = RGBColor(0x10, 0xB9, 0x81)   # Green (#10B981)
BRAND_AMBER   = RGBColor(0xF5, 0x9E, 0x0B)   # Amber (#F59E0B)
BRAND_RED     = RGBColor(0xEF, 0x44, 0x44)   # Red (#EF4444)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY    = RGBColor(0x94, 0xA3, 0xB8)   # Slate 400
DARK_TEXT      = RGBColor(0x1E, 0x29, 0x3B)   # Slate 800
MEDIUM_TEXT    = RGBColor(0x47, 0x55, 0x69)   # Slate 600
SLIDE_BG       = RGBColor(0x0F, 0x17, 0x2A)   # Dark navy
CARD_BG        = RGBColor(0x1E, 0x29, 0x3B)   # Card dark


def set_slide_bg(slide, color=SLIDE_BG):
    """Set slide background to solid color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_accent_bar(slide, left=0, top=0, width=Inches(10), height=Inches(0.06), color=BRAND_BLUE):
    """Add a thin accent bar."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    """Add a simple text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_slide_content(slide, title, subtitle, bullets, y_start=Inches(1.8)):
    """Add title + subtitle + bullet list to a slide."""
    set_slide_bg(slide)
    add_accent_bar(slide, top=Inches(0))

    # Title
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8.4), Inches(0.6),
                 title, font_size=28, color=WHITE, bold=True)

    # Subtitle
    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(1.05), Inches(8.4), Inches(0.5),
                     subtitle, font_size=14, color=LIGHT_GRAY)

    # Bullets
    txBox = slide.shapes.add_textbox(Inches(0.8), y_start, Inches(8.4), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(0xCB, 0xD5, 0xE1)  # Slate 300
        p.font.name = 'Calibri'
        p.space_after = Pt(8)
        p.level = 0


def add_two_column_slide(slide, title, left_title, left_items, right_title, right_items):
    """Add a two-column layout slide."""
    set_slide_bg(slide)
    add_accent_bar(slide, top=Inches(0))

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8.4), Inches(0.6),
                 title, font_size=28, color=WHITE, bold=True)

    # Left column header
    add_text_box(slide, Inches(0.8), Inches(1.3), Inches(4), Inches(0.4),
                 left_title, font_size=16, color=BRAND_BLUE, bold=True)

    # Left items
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(4), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0xCB, 0xD5, 0xE1)
        p.font.name = 'Calibri'
        p.space_after = Pt(6)

    # Right column header
    add_text_box(slide, Inches(5.3), Inches(1.3), Inches(4), Inches(0.4),
                 right_title, font_size=16, color=BRAND_TEAL, bold=True)

    # Right items
    txBox2 = slide.shapes.add_textbox(Inches(5.3), Inches(1.8), Inches(4), Inches(4.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, item in enumerate(right_items):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0xCB, 0xD5, 0xE1)
        p.font.name = 'Calibri'
        p.space_after = Pt(6)


def add_kpi_card(slide, left, top, width, height, label, value, color=BRAND_BLUE):
    """Add a KPI metric card."""
    # Card background
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.fill.background()

    # Color accent bar at top
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.04))
    accent.fill.solid()
    accent.fill.fore_color.rgb = color
    accent.line.fill.background()

    # Value
    add_text_box(slide, left + Inches(0.15), top + Inches(0.15), width - Inches(0.3), Inches(0.4),
                 value, font_size=22, color=WHITE, bold=True)

    # Label
    add_text_box(slide, left + Inches(0.15), top + Inches(0.55), width - Inches(0.3), Inches(0.3),
                 label, font_size=11, color=LIGHT_GRAY)


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # Blank layout

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 1 — Title Slide
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)

    # Decorative circle (top-left glow)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(-2), Inches(6), Inches(6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0x1E, 0x30, 0x5A)
    circle.line.fill.background()

    # Decorative circle (bottom-right glow)
    circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7), Inches(4), Inches(5), Inches(5))
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = RGBColor(0x1A, 0x20, 0x45)
    circle2.line.fill.background()

    # Logo badge
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.0), Inches(1.2), Inches(2), Inches(0.9))
    badge.fill.solid()
    badge.fill.fore_color.rgb = BRAND_BLUE
    badge.line.fill.background()
    add_text_box(slide, Inches(4.0), Inches(1.25), Inches(2), Inches(0.8),
                 "ET", font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Title
    add_text_box(slide, Inches(1), Inches(2.4), Inches(8), Inches(1),
                 "EduTrack", font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Subtitle
    add_text_box(slide, Inches(1), Inches(3.3), Inches(8), Inches(0.6),
                 "School Management Redefined", font_size=24, color=BRAND_BLUE, bold=False, alignment=PP_ALIGN.CENTER)

    # Description
    add_text_box(slide, Inches(1.5), Inches(4.2), Inches(7), Inches(0.8),
                 "A Complete Multi-Tenant Cloud SaaS Platform for Modern Schools",
                 font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # Company
    add_text_box(slide, Inches(1), Inches(6.0), Inches(8), Inches(0.4),
                 "Covenant Synergy Private Limited", font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    add_accent_bar(slide, top=Inches(5.6), left=Inches(3.5), width=Inches(3), color=BRAND_BLUE)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 2 — About EduTrack
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bullet_slide_content(slide, "About EduTrack",
        "What is EduTrack and why does it matter?",
        [
            "☁️  Cloud-based School Management System built as a multi-tenant SaaS platform",
            "🏫  Target audience: K-12 schools, colleges, coaching institutes, and education groups",
            "🌐  Multi-tenant architecture — one platform serves unlimited schools",
            "🎨  Per-school custom branding with dedicated subdomains (school1.edutrack.com)",
            "📊  15+ fully functional modules covering all school operations",
            "🗄️  30+ database tables with complete relational data models",
            "🔗  100+ REST API endpoints powering every feature",
            "📱  30+ responsive frontend screens (desktop + mobile optimized)",
            "👥  6 distinct user roles with fine-grained access control",
            "🚀  Deployed on Vercel (serverless) with AWS RDS PostgreSQL",
        ])

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 3 — Problems Solved
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_two_column_slide(slide, "Problems Solved by EduTrack",
        "❌ Traditional Challenges",
        [
            "• Paper-based student records lost easily",
            "• Manual fee collection with no receipts",
            "• Attendance taken on registers, never analyzed",
            "• No visibility for parents on child's progress",
            "• Timetables made on paper, full of conflicts",
            "• Exam marks calculated manually, error-prone",
            "• No centralized school data analytics",
            "• Each school needs separate software installation",
        ],
        "✅ EduTrack Solutions",
        [
            "• Digital student profiles with instant search",
            "• Automated invoicing with PDF receipts",
            "• Real-time digital roll call with analytics",
            "• Parent portal for grades, fees & attendance",
            "• Automated conflict-free timetable scheduling",
            "• Digital marks entry with grade reports",
            "• Unified dashboard with real-time KPIs",
            "• One SaaS platform serves all schools",
        ])

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 4 — System Architecture
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_accent_bar(slide, top=Inches(0))
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8.4), Inches(0.6),
                 "System Architecture", font_size=28, color=WHITE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(1.05), Inches(8.4), Inches(0.4),
                 "Modern serverless cloud architecture with multi-tenant data isolation", font_size=14, color=LIGHT_GRAY)

    # Architecture boxes
    layers = [
        ("Frontend", "Next.js 14 · React 18 · TailwindCSS\nFramer Motion · Lucide Icons · SWR", BRAND_BLUE, Inches(0.8)),
        ("Backend API", "NestJS 10 · TypeScript · Prisma ORM\nPassport.js · JWT · bcrypt", BRAND_INDIGO, Inches(3.4)),
        ("Database", "PostgreSQL (AWS RDS)\n30+ tables · Multi-tenant isolation\nConnection pooling", BRAND_PURPLE, Inches(6.0)),
    ]

    for label, desc, color, top_pos in layers:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top_pos, Inches(3.5), Inches(1.0))
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.color.rgb = color
        shape.line.width = Pt(2)

        add_text_box(slide, Inches(1.0), top_pos + Inches(0.05), Inches(3.2), Inches(0.3),
                     label, font_size=16, color=color, bold=True)
        add_text_box(slide, Inches(1.0), top_pos + Inches(0.35), Inches(3.2), Inches(0.6),
                     desc, font_size=11, color=LIGHT_GRAY)

    # Right side - services
    services = [
        ("☁️  Vercel Serverless", "Zero-config deployment\nEdge + Serverless Functions", Inches(1.5)),
        ("🔐  Authentication", "JWT tokens + OTP\nRole-based guards (RBAC)", Inches(3.0)),
        ("🏢  Multi-Tenant", "Subdomain resolution\nAsyncLocalStorage context\nPer-tenant data isolation", Inches(4.5)),
        ("📦  AWS Services", "S3 for file storage\nRDS for managed PostgreSQL", Inches(6.2)),
    ]

    for label, desc, top_pos in services:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.3), top_pos, Inches(4.0), Inches(1.0))
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.fill.background()

        add_text_box(slide, Inches(5.5), top_pos + Inches(0.05), Inches(3.6), Inches(0.3),
                     label, font_size=14, color=WHITE, bold=True)
        add_text_box(slide, Inches(5.5), top_pos + Inches(0.35), Inches(3.6), Inches(0.6),
                     desc, font_size=10, color=LIGHT_GRAY)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 5 — User Roles
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_accent_bar(slide, top=Inches(0))
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8.4), Inches(0.6),
                 "User Roles & Access Control", font_size=28, color=WHITE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(1.0), Inches(8.4), Inches(0.4),
                 "6 distinct roles with fine-grained, JWT-enforced permissions", font_size=14, color=LIGHT_GRAY)

    roles = [
        ("SUPER ADMIN", "Platform-level access. Manage all tenants and global settings.", BRAND_RED),
        ("SCHOOL ADMIN", "Full school access. All modules, settings, users, and reports.", BRAND_BLUE),
        ("TEACHER", "Attendance entry, marks entry, complaint submission, schedule.", BRAND_INDIGO),
        ("STAFF", "Billing, library, expenses, student management operations.", BRAND_PURPLE),
        ("STUDENT", "View grades, attendance, schedule, library, and own profile.", BRAND_TEAL),
        ("PARENT", "Monitor child's grades, attendance, and fee statements.", BRAND_AMBER),
    ]

    for i, (role, desc, color) in enumerate(roles):
        col = i % 3
        row = i // 3
        left = Inches(0.8) + col * Inches(3.1)
        top = Inches(1.6) + row * Inches(2.5)

        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.8), Inches(2.0))
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.fill.background()

        # Color accent
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(2.8), Inches(0.04))
        accent.fill.solid()
        accent.fill.fore_color.rgb = color
        accent.line.fill.background()

        add_text_box(slide, left + Inches(0.15), top + Inches(0.2), Inches(2.5), Inches(0.35),
                     role, font_size=14, color=color, bold=True)
        add_text_box(slide, left + Inches(0.15), top + Inches(0.65), Inches(2.5), Inches(1.2),
                     desc, font_size=11, color=LIGHT_GRAY)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 6 — School Registration & Onboarding
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bullet_slide_content(slide, "School Registration & Onboarding",
        "Self-service registration — schools go live in minutes",
        [
            "📝  Multi-step registration wizard collects school details",
            "🏫  School name, type, admin name, mobile, email, address, academic year",
            "🌐  Automatic subdomain allocation (e.g., greenfield.edutrack.com)",
            "👤  Admin user auto-created with SCHOOL_ADMIN role and JWT session",
            "📋  Post-registration setup checklist tracks profile completion %",
            "🎨  School branding: custom logo, name, subtitle on every screen",
            "🏦  Banking & UPI setup: Bank name, IFSC, Account, GPay, PhonePe, UPI QR",
            "✅  13-field completion tracker ensures schools are fully configured",
            "🔄  Settings sync across Tenant, SchoolSetup, and User tables",
        ])

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 7 — Dashboard & Analytics
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_accent_bar(slide, top=Inches(0))
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8.4), Inches(0.6),
                 "Dashboard & Real-Time Analytics", font_size=28, color=WHITE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(1.0), Inches(8.4), Inches(0.4),
                 "Bird's-eye view of school operations with live KPIs and trend analysis", font_size=14, color=LIGHT_GRAY)

    # KPI Cards
    kpis = [
        ("Total Students", "450+", BRAND_BLUE),
        ("Teachers & Staff", "35+", BRAND_INDIGO),
        ("Total Revenue", "₹12.5L", BRAND_GREEN),
        ("Attendance Rate", "92.3%", BRAND_TEAL),
        ("Academic Avg", "78.5%", BRAND_PURPLE),
    ]

    for i, (label, value, color) in enumerate(kpis):
        add_kpi_card(slide, Inches(0.8) + i * Inches(1.75), Inches(1.6), Inches(1.55), Inches(1.0), label, value, color)

    # Features list
    features = [
        "📈  6-month rolling financial chart: Fee Collection vs Salary Expense vs Net Revenue",
        "📊  Month-over-month trend arrows for all KPIs (students, revenue, attendance, academic)",
        "👤  Recent Admissions feed — latest 10 enrolled students with class & joining date",
        "💰  Recent Payments feed — latest fee payments and salary disbursements combined",
        "🔄  Cached queries with 30-second TTL for optimal multi-user performance",
        "📱  Fully responsive — KPI cards stack vertically on mobile screens",
    ]

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(3.0), Inches(8.4), Inches(4.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, feat in enumerate(features):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = feat
        p.font.size = Pt(13)
        p.font.color.rgb = RGBColor(0xCB, 0xD5, 0xE1)
        p.font.name = 'Calibri'
        p.space_after = Pt(6)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 8 — Student Management
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bullet_slide_content(slide, "Student Management",
        "Centralized student directory with full profile lifecycle",
        [
            "🔍  Searchable directory with class/section filtering across 1000+ records",
            "👤  Student profile cards: Photo, name, roll number, class, section, guardian details",
            "📝  Edit student profile with inline modal (name, phone, email, guardian info)",
            "🔢  Automatic roll number generation during admission",
            "💰  Financial status badge: payment progress percentage per student",
            "📤  Bulk student import via structured data upload",
            "🗑️  Bulk delete with actor-audit logging for compliance",
            "👨‍👩‍👧  Parent linkage: Father name, Mother name, Aadhar number tracking",
            "📊  Class and Section read-only display in edit mode (change via promotions)",
        ])

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 9 — Admissions Wizard
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_accent_bar(slide, top=Inches(0))
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8.4), Inches(0.6),
                 "Admissions Wizard", font_size=28, color=WHITE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(1.0), Inches(8.4), Inches(0.4),
                 "End-to-end admission from data entry to fee structure in one workflow", font_size=14, color=LIGHT_GRAY)

    steps = [
        ("Step 1", "Student Details", "Name, phone, email,\nfather/mother name,\nAadhar number", BRAND_BLUE),
        ("Step 2", "Class & Section", "Select academic year,\nclass, and section\nfor enrollment", BRAND_INDIGO),
        ("Step 3", "Fee Products", "Auto-load PriceBook\nfor selected class.\nSelect applicable fees", BRAND_PURPLE),
        ("Step 4", "Concessions", "Apply concession\namount. Calculate\nfinal fee structure", BRAND_TEAL),
        ("Step 5", "Confirm & Create", "Auto-create User,\nStudentProfile, Opportunity,\nand fee line items", BRAND_GREEN),
    ]

    for i, (step, title, desc, color) in enumerate(steps):
        left = Inches(0.3) + i * Inches(1.9)
        top = Inches(1.8)

        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(1.7), Inches(3.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.fill.background()

        # Step number badge
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.55), top + Inches(0.15), Inches(0.6), Inches(0.6))
        badge.fill.solid()
        badge.fill.fore_color.rgb = color
        badge.line.fill.background()
        add_text_box(slide, left + Inches(0.55), top + Inches(0.2), Inches(0.6), Inches(0.5),
                     str(i + 1), font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

        add_text_box(slide, left + Inches(0.1), top + Inches(0.9), Inches(1.5), Inches(0.35),
                     title, font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, left + Inches(0.1), top + Inches(1.4), Inches(1.5), Inches(1.8),
                     desc, font_size=10, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # Arrow connectors text
    for i in range(4):
        left = Inches(2.0) + i * Inches(1.9)
        add_text_box(slide, left, Inches(3.2), Inches(0.3), Inches(0.3),
                     "→", font_size=20, color=BRAND_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

    # Bottom note
    add_text_box(slide, Inches(0.8), Inches(5.8), Inches(8.4), Inches(0.5),
                 "⚡ Single-click enrollment: Student profile, user account, fee structure, and opportunity all created atomically in one database transaction",
                 font_size=11, color=BRAND_TEAL)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 10 — Fee Management & Billing
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_two_column_slide(slide, "Fee Management & Billing",
        "💳 Fee Collection",
        [
            "• Search students by name/class/section",
            "• View unpaid fees per student opportunity",
            "• Multi-item invoice creation",
            "• Payment: Cash / UPI / Bank Transfer / Card",
            "• Bank details attachment to invoices",
            "• Individual & bulk discount application",
            "• Invoice voiding and reversal",
            "• Professional PDF receipt generation",
        ],
        "📋 PriceBook Configuration",
        [
            "• Create fee products: Tuition, Lab, Sports, etc.",
            "• Class-specific pricing per academic year",
            "• Activate/deactivate individual fee items",
            "• Salesforce-style Opportunity/OLI model",
            "• Automatic fee assignment on admission",
            "• Bulk student import with fee structure",
            "• Revenue tracking and collection reports",
            "• Outstanding receivables monitoring",
        ])

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 11 — Invoice PDF Receipt
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_accent_bar(slide, top=Inches(0))
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8.4), Inches(0.6),
                 "Professional Fee Receipt / Invoice", font_size=28, color=WHITE, bold=True)

    # Mock invoice card
    inv_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.3), Inches(7), Inches(5.5))
    inv_shape.fill.solid()
    inv_shape.fill.fore_color.rgb = WHITE
    inv_shape.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)

    inv_lines = [
        ("School Logo & Name", Inches(1.5), Inches(1.5), 16, DARK_TEXT, True),
        ("Invoice #INV-2026-0142                                                  Date: July 04, 2026", Inches(1.7), Inches(2.1), 11, MEDIUM_TEXT, False),
        ("Student: Rahul Sharma  |  Class: Grade 10 - Section A  |  Roll No: 1042", Inches(1.7), Inches(2.5), 11, DARK_TEXT, False),
        ("─────────────────────────────────────────────────────", Inches(1.7), Inches(2.9), 10, LIGHT_GRAY, False),
        ("Fee Item                                  Amount", Inches(1.7), Inches(3.2), 12, DARK_TEXT, True),
        ("Tuition Fee                               ₹ 25,000.00", Inches(1.7), Inches(3.6), 11, MEDIUM_TEXT, False),
        ("Computer Lab Fee                      ₹   3,500.00", Inches(1.7), Inches(3.9), 11, MEDIUM_TEXT, False),
        ("Sports Fee                                  ₹   2,000.00", Inches(1.7), Inches(4.2), 11, MEDIUM_TEXT, False),
        ("Library Fee                                  ₹   1,500.00", Inches(1.7), Inches(4.5), 11, MEDIUM_TEXT, False),
        ("Discount (10%)                           – ₹  3,200.00", Inches(1.7), Inches(4.8), 11, BRAND_RED, False),
        ("─────────────────────────────────────────────────────", Inches(1.7), Inches(5.1), 10, LIGHT_GRAY, False),
        ("TOTAL PAID                               ₹ 28,800.00", Inches(1.7), Inches(5.4), 14, BRAND_BLUE, True),
        ("Payment: UPI  |  Bank: SBI, Main Branch, IFSC: SBIN0001234", Inches(1.7), Inches(5.9), 10, MEDIUM_TEXT, False),
        ("Print-ready • Professional branding • Digital PDF download", Inches(1.7), Inches(6.3), 10, BRAND_TEAL, False),
    ]

    for text, left, top, size, color, bold in inv_lines:
        add_text_box(slide, Inches(left - 0.2) if isinstance(left, float) else left,
                     Inches(top - 0.2) if isinstance(top, float) else top,
                     Inches(6.5), Inches(0.3), text, font_size=size, color=color, bold=bold)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 12 — Attendance Tracking
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bullet_slide_content(slide, "Attendance Tracking",
        "Paperless digital roll call with real-time analytics",
        [
            "📅  Class / Section / Date picker for targeted attendance entry",
            "☑️  Interactive checkbox-based student roll call (Present / Absent / Late / Excused)",
            "📊  Session summary: present count, absent count, total students, attendance %",
            "📆  Calendar-based historical attendance grid with color-coded days",
            "📈  Class-wide attendance report with monthly/weekly trend charts",
            "👤  Individual student attendance calendar heatmap and summary stats",
            "📋  Daily and monthly aggregate summaries across all classes",
            "👩‍🏫  Teacher attribution: every session linked to the staff member who took it",
            "🔄  Edit and update existing attendance sessions",
            "🗑️  Delete incorrect sessions with proper cleanup",
        ])

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 13 — Timetable & Scheduling
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_two_column_slide(slide, "Timetable & Workload Scheduler",
        "📅 Scheduling Engine",
        [
            "• Weekly grid builder (Mon – Sat)",
            "• Subject-Teacher-Period assignment",
            "• Automated conflict detection & prevention",
            "• Period timing configuration (start/end)",
            "• Substitute teacher assignment per period",
            "• Bulk subject import",
            "• Bulk teacher import with skill mapping",
            "• Class-Section creation and management",
        ],
        "📊 Workload Analytics",
        [
            "• Teacher workload summary across all classes",
            "• Class workload summary across all subjects",
            "• Individual teacher drill-down schedule",
            "• Gap period identification per teacher",
            "• Leaser period analysis",
            "• Periods-per-week tracking per assignment",
            "• Teacher skill levels: Beginner/Intermediate/Expert",
            "• Years of experience per subject skill",
        ])

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 14 — Examinations & Grading
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bullet_slide_content(slide, "Examinations & Grading",
        "Complete exam lifecycle from type creation to grade reports",
        [
            "📝  Custom exam type management: Unit Test, Mid-Term, Final, etc. (CRUD)",
            "🗓️  Create exams with name, type, class-section, and date",
            "📊  Subject-wise marks entry grid for selected class-section",
            "✏️  Per-student marks input (0-100) with optional remarks",
            "💾  Batch save: submit all marks for a class in one click",
            "📈  Class-section grade report: aggregate performance across subjects",
            "📉  Mark distribution curve: Failed / Below Average / Average / First Division / High Distinction",
            "📋  Pass rate analytics: automated pass/fail ratio calculation",
            "🏆  Subject-wise performance breakdown for targeted intervention",
        ])

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 15 — Student Promotions
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bullet_slide_content(slide, "Student Promotions",
        "Bulk promotion from one academic year/class to the next",
        [
            "📆  Select source academic year and target academic year",
            "🏫  Filter promotion candidates by class and section",
            "☑️  Bulk student selection with select-all / individual checkboxes",
            "🔄  Automatic class-section reassignment upon promotion",
            "💰  New invoice/opportunity generation for the promoted year",
            "⚡  Entire promotion workflow in a single atomic database transaction",
            "📊  Eliminates manual re-enrollment — promotes entire classes in seconds",
        ])

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 16 — Teacher & Staff Management
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_two_column_slide(slide, "Teacher & Staff Management",
        "👩‍🏫 Profile & Skills",
        [
            "• Teacher profiles: designation, qualification, employee ID",
            "• Subject skill mapping (Beginner/Intermediate/Expert)",
            "• Years of experience per subject",
            "• Class-subject assignments with periods/week",
            "• Teacher schedule and timetable view",
            "• Teacher behavior cases history",
            "• Teaching staff vs. non-teaching staff separation",
        ],
        "💰 Payroll & HR",
        [
            "• Salary structure: Basic + Allowances",
            "• Deductions: standard deductions + PF",
            "• Net salary auto-calculation",
            "• Individual salary payment per month",
            "• Bulk salary payment for all staff",
            "• Salary invoice/receipt generation",
            "• Payment history tracking",
            "• Joining date and status (Active/Suspended/Resigned)",
        ])

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 17 — Expense Management
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bullet_slide_content(slide, "Expense Management",
        "Track, categorize, and approve all school expenditures",
        [
            "📝  Create expenses: amount, category, date, description, payment mode",
            "📂  Category-based classification and filtering",
            "🔄  Status workflow: Pending → Approved → Paid / Rejected",
            "📊  Expense summary: total expenses, category breakdown",
            "✏️  Edit and update expense records",
            "🗑️  Delete expense entries",
            "💰  Integrated with dashboard financial analytics",
            "📈  Month-over-month expense trend tracking",
        ])

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 18 — Library Management
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bullet_slide_content(slide, "Library Management",
        "Digital library with barcode-based book tracking and fine management",
        [
            "📚  Book registration: Title, Author, ISBN, Category, Total Copies",
            "🏷️  BookCopy tracking with unique barcodes",
            "📤  Barcode-based book issue to students and staff",
            "📥  Book return with automatic due date checking",
            "💰  Late fee / fine calculation on overdue books",
            "📊  Copy status monitoring: Available / Issued / Lost / Damaged",
            "📋  Complete borrow log history with borrower details",
            "🔍  Available copies count auto-updated on issue/return",
        ])

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 19 — Complaint Box
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bullet_slide_content(slide, "Complaint Box & Behavior Tracking",
        "Structured incident management for student behavior",
        [
            "📝  Submit behavior records: Complaints (negative) and Praise (positive)",
            "📂  Category and priority classification (High / Medium)",
            "🔄  Status workflow: New → In Progress → Resolved → Closed",
            "👩‍🏫  Teacher attribution: every case linked to the reporting teacher",
            "🔍  Class-based student search for incident reporting",
            "📊  Student behavior statistics and incident history",
            "📅  Academic year filtering for period-specific views",
            "📋  Pending cases dashboard for admin review",
            "🔗  Linked to teacher profile: view all cases reported by a teacher",
        ])

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 20 — Reports & Analytics Hub
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_accent_bar(slide, top=Inches(0))
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8.4), Inches(0.6),
                 "Reports & Analytics Hub", font_size=28, color=WHITE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(1.0), Inches(8.4), Inches(0.4),
                 "Consolidated reporting across all modules for data-driven decisions", font_size=14, color=LIGHT_GRAY)

    reports = [
        ("📊 Enrollment Demographics", "Student counts by class\nAdmission timeline charts\nEnrollment growth trends", BRAND_BLUE),
        ("💰 Financial Statements", "Total revenue & receivables\nExpense breakdown\nNet cashflow analysis", BRAND_GREEN),
        ("📈 Grading Analytics", "Mark distribution curve\nPass/fail rate analysis\nSubject-wise performance", BRAND_PURPLE),
        ("📋 Cashflow Reports", "Transaction-level detail\nFee revenue breakdown\nSalary disbursements", BRAND_TEAL),
    ]

    for i, (title, desc, color) in enumerate(reports):
        col = i % 2
        row = i // 2
        left = Inches(0.8) + col * Inches(4.5)
        top = Inches(1.7) + row * Inches(2.5)

        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(4.0), Inches(2.1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.fill.background()

        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.06), Inches(2.1))
        accent.fill.solid()
        accent.fill.fore_color.rgb = color
        accent.line.fill.background()

        add_text_box(slide, left + Inches(0.25), top + Inches(0.2), Inches(3.5), Inches(0.35),
                     title, font_size=15, color=WHITE, bold=True)
        add_text_box(slide, left + Inches(0.25), top + Inches(0.65), Inches(3.5), Inches(1.3),
                     desc, font_size=12, color=LIGHT_GRAY)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 21 — Multi-Tenant Architecture
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bullet_slide_content(slide, "Multi-Tenant SaaS Architecture",
        "One platform, unlimited schools — with complete data isolation",
        [
            "🌐  Subdomain-based tenant resolution: school1.edutrack.com → Tenant ID",
            "📨  X-Tenant-ID header for API requests from SPA clients",
            "❓  Query parameter fallback: ?tenant=school1 for flexibility",
            "🧵  AsyncLocalStorage-based tenant context propagation (Node.js)",
            "🔒  Per-tenant data filtering on ALL 30+ database tables automatically",
            "🎨  School-specific branding: custom logo, name, subtitle on every screen",
            "🏦  Per-school banking & UPI configuration for fee receipts",
            "🔐  Subdomain uniqueness validation prevents conflicts",
            "📊  Data isolation ensures School A can NEVER see School B's data",
            "⚡  Single codebase deployment serves all schools simultaneously",
        ])

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 22 — Security & Compliance
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_two_column_slide(slide, "Security & Compliance",
        "🔐 Authentication",
        [
            "• JWT token-based session management",
            "• bcrypt password hashing (10 salt rounds)",
            "• OTP-based passwordless mobile login",
            "• Token expiry: 24h access, 7d refresh",
            "• Secure cookie-less API authentication",
        ],
        "🛡️ Authorization & Audit",
        [
            "• Role-Based Access Control (RBAC)",
            "• NestJS route guards per endpoint",
            "• Tenant isolation at middleware level",
            "• Activity audit logging for all actions",
            "• Entity-level change tracking with JSON diffs",
            "• Connection pooling for serverless security",
        ])

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 23 — Technology Stack
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    add_accent_bar(slide, top=Inches(0))
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8.4), Inches(0.6),
                 "Technology Stack", font_size=28, color=WHITE, bold=True)

    stack = [
        ("Frontend", "Next.js 14 (App Router) · React 18 · TypeScript\nTailwindCSS 3.4 · Framer Motion · Lucide React\nSWR · Axios · Responsive Design", BRAND_BLUE),
        ("Backend", "NestJS 10 · TypeScript · Prisma 5 ORM\nPassport.js · JWT · bcrypt · class-validator\nclass-transformer · CSV Parser", BRAND_INDIGO),
        ("Database", "PostgreSQL (AWS RDS) · 30+ Prisma models\nUUID primary keys · Multi-tenant indexed\nRelational schema with cascade deletes", BRAND_PURPLE),
        ("Infrastructure", "Vercel Serverless (Frontend + Backend)\nAWS S3 (File Storage) · AWS RDS (Database)\nGitHub CI/CD · Zero-downtime deployments", BRAND_TEAL),
    ]

    for i, (category, techs, color) in enumerate(stack):
        top = Inches(1.3) + i * Inches(1.5)
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top, Inches(8.4), Inches(1.3))
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.fill.background()

        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), top, Inches(0.06), Inches(1.3))
        accent.fill.solid()
        accent.fill.fore_color.rgb = color
        accent.line.fill.background()

        add_text_box(slide, Inches(1.1), top + Inches(0.1), Inches(2.0), Inches(0.3),
                     category, font_size=16, color=color, bold=True)
        add_text_box(slide, Inches(3.2), top + Inches(0.1), Inches(5.8), Inches(1.0),
                     techs, font_size=12, color=LIGHT_GRAY)

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 24 — Future Enhancements
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    add_bullet_slide_content(slide, "Future Enhancements",
        "Planned modules and features for the next release cycle",
        [
            "🔮  Parent Portal — Online grade viewing, attendance tracking, UPI fee payments",
            "📚  Teacher Portal — Homework tracking, lesson planner, syllabus progress",
            "🚌  Transport Management — Route mapping, vehicle tracking, driver allocation",
            "🏠  Hostel Management — Room allocation, occupancy tracking, meal plans",
            "🔔  Real-Time Notifications — WebSocket push notifications across all roles",
            "📧  SMS/Email Integration — Twilio/SendGrid for automated communication",
            "📱  Mobile App — React Native companion app for parents & teachers",
            "🤖  AI Analytics — Predictive student performance & dropout risk analysis",
            "💳  Payment Gateway — Razorpay/Stripe integration for online fee payments",
            "📄  Advanced Report Cards — Printable PDF report cards with school branding",
        ])

    # ═══════════════════════════════════════════════════════════════════
    # SLIDE 25 — Conclusion & CTA
    # ═══════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)

    # Decorative circles
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1), Inches(-1), Inches(4), Inches(4))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0x1E, 0x30, 0x5A)
    circle.line.fill.background()

    circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7), Inches(5), Inches(4), Inches(4))
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = RGBColor(0x1A, 0x20, 0x45)
    circle2.line.fill.background()

    add_text_box(slide, Inches(1), Inches(0.8), Inches(8), Inches(0.8),
                 "EduTrack — By the Numbers", font_size=32, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_accent_bar(slide, top=Inches(1.6), left=Inches(3.5), width=Inches(3), color=BRAND_BLUE)

    stats = [
        ("15+", "Modules", BRAND_BLUE),
        ("30+", "DB Tables", BRAND_INDIGO),
        ("100+", "APIs", BRAND_PURPLE),
        ("30+", "Screens", BRAND_TEAL),
        ("6", "User Roles", BRAND_GREEN),
    ]

    for i, (val, label, color) in enumerate(stats):
        left = Inches(0.5) + i * Inches(1.85)
        add_kpi_card(slide, left, Inches(2.2), Inches(1.65), Inches(1.1), label, val, color)

    add_text_box(slide, Inches(1), Inches(4.0), Inches(8), Inches(0.7),
                 "Ready to Transform Your School Management?",
                 font_size=24, color=BRAND_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(1), Inches(4.8), Inches(8), Inches(0.5),
                 "Multi-Tenant SaaS  •  Cloud-Native  •  Zero Infrastructure Hassle",
                 font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # CTA box
    cta = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), Inches(5.5), Inches(5), Inches(1.2))
    cta.fill.solid()
    cta.fill.fore_color.rgb = BRAND_BLUE
    cta.line.fill.background()

    add_text_box(slide, Inches(2.5), Inches(5.55), Inches(5), Inches(0.5),
                 "Request a Demo Today", font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(2.5), Inches(6.0), Inches(5), Inches(0.4),
                 "contact@covenantsynergy.com  |  www.edutrack.app", font_size=12, color=RGBColor(0xBF, 0xDB, 0xFE), alignment=PP_ALIGN.CENTER)

    # Company footer
    add_text_box(slide, Inches(1), Inches(7.0), Inches(8), Inches(0.3),
                 "© 2026 Covenant Synergy Private Limited. All rights reserved.",
                 font_size=10, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # ═══════════════════════════════════════════════════════════════════
    # Save
    # ═══════════════════════════════════════════════════════════════════
    output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'EduTrack_Presentation_Generated.pptx')
    prs.save(output_path)
    print(f"\n[OK] Presentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")
    print(f"Brand colors: Navy dark background with blue/indigo/purple accent palette")
    print(f"\nOpen the file in Microsoft PowerPoint or Google Slides to view.")


if __name__ == '__main__':
    create_presentation()
